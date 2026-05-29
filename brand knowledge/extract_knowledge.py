import os
import docx2txt
import pptx
import fitz  # PyMuPDF
import json

def extract_text_from_docx(file_path):
    try:
        return docx2txt.process(file_path)
    except Exception as e:
        return f"Error extracting DOCX: {e}"

def extract_text_from_pptx(file_path):
    try:
        prs = pptx.Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    except Exception as e:
        return f"Error extracting PPTX: {e}"

def extract_text_from_pdf(file_path):
    try:
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text()
        return text
    except Exception as e:
        return f"Error extracting PDF: {e}"

def main():
    base_path = r"C:\Users\Aayushi\sleepycat-brand\brand knowledge"
    files = [
        "Brand induction_ Sleepycat 2025.pptx",
        "Brief Doc_ Brand Truth Sleepycat.docx",
        "New Brand Guidelines SleepyCat_Celebrating the Joy of Rest.pptx",
        "Sleepycat R4.pdf"
    ]
    
    knowledge_base = {}
    
    for file_name in files:
        file_path = os.path.join(base_path, file_name)
        print(f"Processing: {file_name}")
        if file_name.endswith(".docx"):
            knowledge_base[file_name] = extract_text_from_docx(file_path)
        elif file_name.endswith(".pptx"):
            knowledge_base[file_name] = extract_text_from_pptx(file_path)
        elif file_name.endswith(".pdf"):
            knowledge_base[file_name] = extract_text_from_pdf(file_path)
            
    with open("brand_knowledge_dump.json", "w", encoding="utf-8") as f:
        json.dump(knowledge_base, f, indent=4)
    print("Extraction complete. Saved to brand_knowledge_dump.json")

if __name__ == "__main__":
    main()
